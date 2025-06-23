# create_files.py

from docx import Document
from pptx import Presentation
import datetime
import os

def create_office_file(file_type: str, topic: str, content: str) -> bool:
    try:
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = topic.replace(" ", "_")

        if file_type == "doc":
            doc = Document()
            doc.add_heading(topic.title(), 0)
            doc.add_paragraph(content)
            path = f"{safe_topic}_{timestamp}.docx"
            doc.save(path)

        elif file_type == "ppt":
            ppt = Presentation()
            slide_layout = ppt.slide_layouts[1]

            for line in content.split('\n'):
                slide = ppt.slides.add_slide(slide_layout)
                slide.shapes.title.text = topic.title()
                slide.placeholders[1].text = line.strip()

            path = f"{safe_topic}_{timestamp}.pptx"
            ppt.save(path)

        else:
            return False

        os.startfile(path)
        return True

    except Exception as e:
        print(f"Error creating file: {e}")
        return False
